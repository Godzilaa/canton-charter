#!/usr/bin/env python3
"""
Generate Charter hackathon pitch deck as a .pptx file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ──────────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x0A, 0x0A, 0x0A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT     = RGBColor(0x00, 0xB4, 0xD8)   # cyan-blue
ACCENT2    = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
DARK_GREY  = RGBColor(0x1A, 0x1A, 0x2E)   # near-black navy
MID_GREY   = RGBColor(0x44, 0x44, 0x55)
LIGHT_GREY = RGBColor(0xE8, 0xE8, 0xF0)
GREEN      = RGBColor(0x2D, 0xC6, 0x53)
ORANGE     = RGBColor(0xFF, 0x8C, 0x00)

FONT = "Courier New"  # brutalist mono — matches the UI

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def blank_slide():
    blank_layout = prs.slide_layouts[6]   # truly blank
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=DARK_GREY):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txb(slide, text, l, t, w, h,
        size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
        wrap=True, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = FONT
    run.font.italic = italic
    return box

def add_line(slide, x1, y1, x2, y2, color=ACCENT, width_pt=2):
    from pptx.util import Pt as _Pt
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = _Pt(width_pt)

def multi_para(slide, paras, l, t, w, h, size=16, color=WHITE, leading_color=ACCENT):
    """Add a textbox with multiple bullet paragraphs."""
    box = slide.shapes.add_textbox(l, t, w, h)
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    for i, (bullet, rest) in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        if bullet:
            r = p.add_run()
            r.text = bullet + " "
            r.font.color.rgb = leading_color
            r.font.size = Pt(size)
            r.font.bold = True
            r.font.name = FONT
        r2 = p.add_run()
        r2.text = rest
        r2.font.color.rgb = color
        r2.font.size = Pt(size)
        r2.font.name = FONT
    return box

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)

# accent bar top
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)
# accent bar bottom
add_rect(s, 0, H - Inches(0.07), W, Inches(0.07), ACCENT)

# big wordmark
txb(s, "CHARTER", Inches(1), Inches(1.6), Inches(11), Inches(2.2),
    size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# tagline
txb(s, "Agentic Payments on Canton Network",
    Inches(1), Inches(3.7), Inches(11), Inches(0.7),
    size=26, bold=False, color=ACCENT2, align=PP_ALIGN.CENTER)

# sub-tagline
txb(s, "DAML Smart Contracts  ×  AI Agents  ×  x402 Protocol",
    Inches(1), Inches(4.4), Inches(11), Inches(0.6),
    size=16, color=MID_GREY, align=PP_ALIGN.CENTER)

# hackathon label
txb(s, "Canton Network Hackathon  //  2026",
    Inches(1), Inches(5.5), Inches(11), Inches(0.5),
    size=13, color=MID_GREY, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "THE PROBLEM", Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
    size=11, color=ACCENT, bold=True)
txb(s, "AI agents need to spend money.\nEnterprises have no safe way to let them.",
    Inches(0.8), Inches(0.75), Inches(11.5), Inches(1.4),
    size=32, bold=True, color=WHITE)

add_line(s, Inches(0.8), Inches(2.2), Inches(12.5), Inches(2.2), ACCENT, 1)

pain_points = [
    ("01", "No policy enforcement — agents can spend beyond intent with no hard stop."),
    ("02", "No audit trail — finance teams can't reconstruct what an agent bought, when, or why."),
    ("03", "No kill switch — revoking an agent's spending authority means revoking credentials."),
    ("04", "API paywalls block agents — most data/compute APIs charge per-call; agents can't self-pay."),
]
for i, (num, text) in enumerate(pain_points):
    y = Inches(2.4) + i * Inches(1.1)
    add_rect(s, Inches(0.8), y, Inches(0.55), Inches(0.55), ACCENT)
    txb(s, num, Inches(0.82), y + Inches(0.04), Inches(0.5), Inches(0.5),
        size=14, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txb(s, text, Inches(1.55), y + Inches(0.05), Inches(10.9), Inches(0.9),
        size=17, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "THE SOLUTION", Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
    size=11, color=ACCENT, bold=True)
txb(s, "Charter",
    Inches(0.8), Inches(0.75), Inches(6), Inches(0.9),
    size=44, bold=True, color=WHITE)
txb(s, "Encode enterprise spending policy as DAML smart contracts.\nAI agents get authority to act — and hard boundaries they cannot cross.",
    Inches(0.8), Inches(1.7), Inches(11.5), Inches(1.1),
    size=18, color=ACCENT2)

# 3 pillars
pillars = [
    (ACCENT,   "POLICY AS CODE",   "Spending limits, categories, per-tx caps, and daily ceilings live in DAML contracts. Agents cannot override them — the ledger enforces."),
    (GREEN,    "x402 PAYMENTS",    "HTTP 402 Payment Required triggers automatic Canton authorization. Agents pay API paywalls machine-to-machine, without human intervention."),
    (ORANGE,   "AUDIT BY DEFAULT", "Every payment creates an immutable PaymentRecord on Canton. Finance gets a tamper-proof ledger — no spreadsheets, no guesswork."),
]
for i, (col, title, body) in enumerate(pillars):
    x = Inches(0.8) + i * Inches(4.15)
    add_rect(s, x, Inches(3.1), Inches(3.9), Inches(0.06), col)
    txb(s, title, x, Inches(3.25), Inches(3.9), Inches(0.55),
        size=13, bold=True, color=col)
    txb(s, body, x, Inches(3.85), Inches(3.8), Inches(2.5),
        size=14, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — How It Works (x402 flow)
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "HOW IT WORKS", Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
    size=11, color=ACCENT, bold=True)
txb(s, "The x402 Payment Flow",
    Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
    size=30, bold=True, color=WHITE)

steps = [
    ("1", ACCENT,  "AGENT REQUESTS",     "AI agent calls a paywall API endpoint for market data, compute, or any paid resource."),
    ("2", ACCENT2, "402 RETURNED",        "API responds HTTP 402 with a machine-readable payment spec — amount, currency, category."),
    ("3", GREEN,   "CANTON AUTHORIZES",   "x402 layer calls RequestAuthorization on the DAML SpendingPolicy. Ledger checks limits — on-chain, atomic."),
    ("4", ACCENT,  "TOKEN SENT",          "Authorization contract ID is sent as X-Payment-Token header. API unlocks and responds 200."),
    ("5", ORANGE,  "SETTLE ON-CHAIN",     "Execute choice fires — PaymentAuthorization consumed, immutable PaymentRecord created with API receipt embedded."),
]

for i, (num, col, title, body) in enumerate(steps):
    y = Inches(1.7) + i * Inches(1.05)
    # number circle
    add_rect(s, Inches(0.8), y, Inches(0.5), Inches(0.5), col)
    txb(s, num, Inches(0.82), y + Inches(0.05), Inches(0.46), Inches(0.45),
        size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    # connector line (except last)
    if i < len(steps) - 1:
        add_rect(s, Inches(1.02), y + Inches(0.5), Inches(0.06), Inches(0.55), MID_GREY)
    txb(s, title + "  ", Inches(1.45), y + Inches(0.06), Inches(2.8), Inches(0.45),
        size=13, bold=True, color=col)
    txb(s, body, Inches(4.0), y + Inches(0.07), Inches(8.5), Inches(0.5),
        size=14, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — DAML Smart Contracts
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "DAML SMART CONTRACTS", Inches(0.8), Inches(0.3), Inches(8), Inches(0.5),
    size=11, color=ACCENT, bold=True)
txb(s, "Three Contracts. One Source of Truth.",
    Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
    size=30, bold=True, color=WHITE)

contracts = [
    ("SpendingPolicy",       ACCENT,  [
        "Owned by enterprise — agent cannot modify",
        "maxPerTx  ·  dailyLimit  ·  allowedCategories",
        "requireApprovalAbove threshold (human-in-loop)",
        "Emergency Deactivate / Reactivate kill switch",
    ]),
    ("PaymentAuthorization", GREEN,   [
        "Created per payment request after policy checks pass",
        "Status: Approved | PendingApproval | Rejected",
        "Pending → CFO approves or rejects via dashboard",
        "Approved → Execute choice settles the payment",
    ]),
    ("PaymentRecord",        ORANGE,  [
        "Immutable — no archival choices",
        "Stores vendor, amount, category, purpose, outcome",
        "Embeds x402Token and API receipt ID",
        "Visible to enterprise, agent, and approver",
    ]),
]

for i, (name, col, bullets) in enumerate(contracts):
    x = Inches(0.8) + i * Inches(4.15)
    add_rect(s, x, Inches(1.65), Inches(3.9), Inches(5.3), RGBColor(0x16, 0x16, 0x26))
    add_rect(s, x, Inches(1.65), Inches(3.9), Inches(0.08), col)
    txb(s, name, x + Inches(0.15), Inches(1.8), Inches(3.6), Inches(0.55),
        size=15, bold=True, color=col)
    for j, b in enumerate(bullets):
        by = Inches(2.45) + j * Inches(0.72)
        txb(s, "▸ " + b, x + Inches(0.15), by, Inches(3.6), Inches(0.65),
            size=13, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture Diagram
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "ARCHITECTURE", Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
    size=11, color=ACCENT, bold=True)
txb(s, "Full-Stack Overview",
    Inches(0.8), Inches(0.7), Inches(11), Inches(0.65),
    size=30, bold=True, color=WHITE)

# Layer boxes
layers = [
    (Inches(0.5),  Inches(1.65), Inches(3.5),  Inches(1.4),  RGBColor(0x10,0x22,0x30), ACCENT,  "NEXT.JS DASHBOARD",
     "Spending Policies  ·  Agent Feed\nApprovals  ·  Audit Log  ·  Kill Switch"),
    (Inches(4.5),  Inches(1.65), Inches(4.3),  Inches(1.4),  RGBColor(0x10,0x30,0x18), GREEN,   "AI AGENT LAYER",
     "Claude-powered  ·  charter_tools\nx402Fetch  ·  Canton SDK calls"),
    (Inches(9.3),  Inches(1.65), Inches(3.5),  Inches(1.4),  RGBColor(0x30,0x20,0x10), ORANGE,  "PAYWALL APIS",
     "Market data  ·  Compute\nAny HTTP API (x402-enabled)"),
    (Inches(2.0),  Inches(4.0),  Inches(9.3),  Inches(2.0),  RGBColor(0x08,0x08,0x1A), ACCENT,  "CANTON NETWORK  (DAML Ledger)",
     "SpendingPolicy  ·  PaymentAuthorization  ·  PaymentRecord\nSub-transaction privacy  ·  Atomic execution  ·  Immutable audit"),
]

for lx, ly, lw, lh, bg_col, border_col, title, body in layers:
    add_rect(s, lx, ly, lw, lh, bg_col)
    # top border
    add_rect(s, lx, ly, lw, Inches(0.055), border_col)
    txb(s, title, lx + Inches(0.12), ly + Inches(0.1), lw - Inches(0.2), Inches(0.45),
        size=12, bold=True, color=border_col)
    txb(s, body,  lx + Inches(0.12), ly + Inches(0.6), lw - Inches(0.2), Inches(0.9),
        size=12, color=LIGHT_GREY)

# arrows between top layers → Canton
for ax in [Inches(2.1), Inches(6.55), Inches(10.4)]:
    add_rect(s, ax, Inches(3.1), Inches(0.06), Inches(0.9), ACCENT)
    txb(s, "▼", ax - Inches(0.04), Inches(3.7), Inches(0.2), Inches(0.3),
        size=13, color=ACCENT, bold=True)

# horizontal arrows between top layers
add_rect(s, Inches(4.05), Inches(2.3), Inches(0.4), Inches(0.06), ACCENT2)
txb(s, "►", Inches(4.3), Inches(2.2), Inches(0.3), Inches(0.3), size=13, color=ACCENT2, bold=True)
add_rect(s, Inches(8.85), Inches(2.3), Inches(0.4), Inches(0.06), ACCENT2)
txb(s, "►", Inches(9.1), Inches(2.2), Inches(0.3), Inches(0.3), size=13, color=ACCENT2, bold=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Live Demo
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "LIVE DEMO", Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
    size=11, color=ACCENT, bold=True)
txb(s, "What We Built",
    Inches(0.8), Inches(0.7), Inches(11), Inches(0.65),
    size=30, bold=True, color=WHITE)

demo_items = [
    (ACCENT,  "Trading Demo",          "node agent/trading-demo.js",
     "AI agent fetches BTCUSD + ETHUSD data from paywall API. Full 5-step x402 flow. Every payment lands on Canton as an immutable PaymentRecord."),
    (GREEN,   "Dashboard",             "Next.js  /  Canton Ledger API",
     "Spending Policies page — live DAML contracts, activate/deactivate toggle. Kill switch blocks all new payment requests in one click."),
    (ORANGE,  "Agent Feed",            "/agent  —  auto-refresh every 10s",
     "Real-time stream of PaymentAuthorization contracts. Status badges: APPROVED / PENDING / REJECTED. x402 settlement confirmed inline."),
    (ACCENT2, "Human-in-the-Loop",     "Approvals page",
     "Payments above the threshold pause as PendingApproval. CFO approves or rejects — the ledger records the decision either way."),
]

for i, (col, title, subtitle, body) in enumerate(demo_items):
    row, col_idx = divmod(i, 2)
    x = Inches(0.8) + col_idx * Inches(6.15)
    y = Inches(1.65) + row * Inches(2.55)
    add_rect(s, x, y, Inches(5.85), Inches(2.3), RGBColor(0x12,0x12,0x22))
    add_rect(s, x, y, Inches(5.85), Inches(0.055), col)
    txb(s, title, x + Inches(0.15), y + Inches(0.12), Inches(5.5), Inches(0.5),
        size=16, bold=True, color=col)
    txb(s, subtitle, x + Inches(0.15), y + Inches(0.65), Inches(5.5), Inches(0.35),
        size=11, color=MID_GREY, italic=True)
    txb(s, body, x + Inches(0.15), y + Inches(1.05), Inches(5.55), Inches(1.1),
        size=13, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Why Canton
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "WHY CANTON", Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
    size=11, color=ACCENT, bold=True)
txb(s, "The Right Ledger for Agentic Payments",
    Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
    size=30, bold=True, color=WHITE)

why_items = [
    (ACCENT,  "Sub-transaction Privacy",
     "In multi-party B2B flows, each participant sees only their part of the transaction. No competitor leaks the trade."),
    (GREEN,   "DAML = Policy as Law",
     "Smart contract logic is checked by the ledger, not by the application. Agents cannot bypass limits — the VM enforces them."),
    (ORANGE,  "Atomic Execution",
     "Authorization and settlement are ledger operations. No race conditions, no double-spends, no half-executed payments."),
    (ACCENT2, "Compliance-Ready Audit",
     "PaymentRecord is immutable and permanent. SOC2 / SOX audit trail is a query, not a reconstruction project."),
]

for i, (col, title, body) in enumerate(why_items):
    row, cidx = divmod(i, 2)
    x = Inches(0.8) + cidx * Inches(6.15)
    y = Inches(1.75) + row * Inches(2.4)
    add_rect(s, x, y, Inches(5.85), Inches(2.1), RGBColor(0x10,0x10,0x20))
    txb(s, "◆ " + title, x + Inches(0.2), y + Inches(0.2), Inches(5.4), Inches(0.6),
        size=16, bold=True, color=col)
    txb(s, body, x + Inches(0.2), y + Inches(0.85), Inches(5.4), Inches(1.1),
        size=15, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Traction / What's Done
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)

txb(s, "TRACTION", Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
    size=11, color=ACCENT, bold=True)
txb(s, "What's Shipped",
    Inches(0.8), Inches(0.7), Inches(11), Inches(0.7),
    size=30, bold=True, color=WHITE)

done = [
    "DAML contracts — SpendingPolicy, PaymentAuthorization, PaymentRecord — deployed on Canton devnet",
    "x402 payment layer — full 5-step flow: 402 intercept → Canton authorize → retry → settle → record",
    "AI Agent — Claude-powered trading agent with charter_tools, runCharterTool integration",
    "Next.js Dashboard — spending policies, agent feed (live 10s poll), approvals, audit log, kill switch",
    "Human-in-the-loop — CFO approval flow for payments above configurable threshold",
    "Mock Trading API — x402-enabled endpoint serving market data (quotes + signals) at $1–$5/call",
    "Live on Canton devnet — real ledger, real contracts, real on-chain PaymentRecords",
]

for i, item in enumerate(done):
    y = Inches(1.65) + i * Inches(0.77)
    add_rect(s, Inches(0.8), y + Inches(0.12), Inches(0.35), Inches(0.35), GREEN)
    txb(s, "✓", Inches(0.82), y + Inches(0.1), Inches(0.32), Inches(0.38),
        size=13, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    txb(s, item, Inches(1.3), y + Inches(0.06), Inches(11.2), Inches(0.62),
        size=15, color=LIGHT_GREY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Vision / Closing
# ═══════════════════════════════════════════════════════════════════════════
s = blank_slide()
bg(s, DARK_GREY)
add_rect(s, 0, 0, W, Inches(0.07), ACCENT)
add_rect(s, 0, H - Inches(0.07), W, Inches(0.07), ACCENT)

txb(s, "VISION", Inches(0.8), Inches(0.3), Inches(11), Inches(0.5),
    size=11, color=ACCENT, bold=True)
txb(s, "Every AI agent,\nevery enterprise,\nzero uncontrolled spend.",
    Inches(0.8), Inches(0.7), Inches(11.5), Inches(2.0),
    size=36, bold=True, color=WHITE)

txb(s, "Charter turns spending policy into a first-class ledger primitive. As AI agents proliferate across enterprise workflows, Canton becomes the trust anchor for every dollar they touch — auditable, atomic, and impossible to bypass.",
    Inches(0.8), Inches(2.9), Inches(11.5), Inches(1.5),
    size=17, color=ACCENT2)

add_line(s, Inches(0.8), Inches(4.55), Inches(12.5), Inches(4.55), MID_GREY, 0.5)

txb(s, "CHARTER", Inches(0.8), Inches(4.75), Inches(5), Inches(0.7),
    size=28, bold=True, color=WHITE)
txb(s, "Canton Network Hackathon  //  2026",
    Inches(0.8), Inches(5.45), Inches(6), Inches(0.4),
    size=13, color=MID_GREY)

txb(s, "github.com/godzilaa/charter",
    Inches(8.5), Inches(4.85), Inches(4.3), Inches(0.45),
    size=13, color=ACCENT, align=PP_ALIGN.RIGHT)
txb(s, "pratikvpatil17@gmail.com",
    Inches(8.5), Inches(5.35), Inches(4.3), Inches(0.45),
    size=13, color=MID_GREY, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), '..', 'Charter_Pitch_Deck.pptx')
prs.save(out)
print(f"Saved: {os.path.abspath(out)}")
